"""PPTX parsing: extract text for pre-filtering + convert slides to PNG images."""

import asyncio
import subprocess
from pathlib import Path

from loguru import logger
from pptx import Presentation

from app.core.config import settings

# Keywords that indicate a case page
CASE_KEYWORDS = [
    "customer", "device", "defect", "fa status", "follow up",
    "fail rate", "defect rate", "defect mode", "defect phenomenon",
    "fab", "assembly", "lot",
]
MIN_KEYWORD_MATCHES = 2


def extract_slide_texts(pptx_path: str | Path) -> list[str]:
    """Extract all text from each slide. Returns list of text per slide."""
    prs = Presentation(str(pptx_path))
    slide_texts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            texts.append(text)
        slide_texts.append("\n".join(texts))
    return slide_texts


def pre_filter_slides(slide_texts: list[str]) -> list[bool]:
    """Check each slide for case-related keywords. Returns list of booleans."""
    results = []
    for text in slide_texts:
        text_lower = text.lower()
        matches = sum(1 for kw in CASE_KEYWORDS if kw in text_lower)
        results.append(matches >= MIN_KEYWORD_MATCHES)
    return results


async def convert_pptx_to_images(
    pptx_path: str | Path, output_dir: str | Path
) -> list[Path]:
    """Convert PPTX to per-slide PNG images using LibreOffice headless.

    Returns list of image paths sorted by slide number.
    """
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    pdf_path = output_dir / pptx_path.with_suffix(".pdf").name

    # Step 1: PPTX → PDF
    cmd_pdf = [
        "libreoffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(pptx_path),
    ]
    process = await asyncio.create_subprocess_exec(
        *cmd_pdf,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
    )
    stdout, stderr = await process.communicate()
    if process.returncode != 0:
        raise RuntimeError(f"PPTX to PDF failed: {stderr.decode()}")

    # Step 2: PDF → per-page PNG using pdftoppm (from poppler-utils)
    # Output pattern: slide-01.png, slide-02.png, ...
    prefix = str(output_dir / "slide")
    cmd_png = [
        "pdftoppm",
        "-png",
        "-r", "200",  # 200 DPI
        str(pdf_path),
        prefix,
    ]
    process = await asyncio.create_subprocess_exec(
        *cmd_png,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
    )
    stdout, stderr = await process.communicate()
    if process.returncode != 0:
        raise RuntimeError(f"PDF to PNG failed: {stderr.decode()}")

    # Clean up intermediate PDF
    pdf_path.unlink(missing_ok=True)

    # Collect output images sorted by name
    images = sorted(output_dir.glob("slide-*.png"))
    if not images:
        raise RuntimeError("No slide images generated")

    logger.info("Converted {} slides to PNG", len(images))
    return images
